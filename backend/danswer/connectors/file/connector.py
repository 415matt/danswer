import csv  # type: ignore
import io
import os
import zipfile
from collections.abc import Iterator
from datetime import datetime
from datetime import timezone
from email.parser import Parser as EmailParser
from pathlib import Path
from typing import Any
from typing import IO

import docx2txt  # type: ignore
import openpyxl  # type: ignore
import pptx  # type: ignore
from bs4 import BeautifulSoup
from sqlalchemy.orm import Session

from danswer.configs.app_configs import INDEX_BATCH_SIZE
from danswer.configs.constants import DocumentSource
from danswer.connectors.cross_connector_utils.file_utils import detect_encoding
from danswer.connectors.cross_connector_utils.file_utils import load_files_from_zip
from danswer.connectors.cross_connector_utils.file_utils import read_file
from danswer.connectors.cross_connector_utils.file_utils import read_pdf_file
from danswer.connectors.cross_connector_utils.miscellaneous_utils import time_str_to_utc
from danswer.connectors.file.utils import check_file_ext_is_valid
from danswer.connectors.file.utils import get_file_ext
from danswer.connectors.interfaces import GenerateDocumentsOutput
from danswer.connectors.interfaces import LoadConnector
from danswer.connectors.models import BasicExpertInfo
from danswer.connectors.models import Document
from danswer.connectors.models import Section
from danswer.db.engine import get_sqlalchemy_engine
from danswer.file_store.file_store import get_default_file_store
from danswer.utils.logger import setup_logger

logger = setup_logger()


def _read_files_and_metadata(
    file_name: str,
    db_session: Session,
) -> Iterator[tuple[str, IO, dict[str, Any]]]:
    """Reads the file into IO, in the case of a zip file, yields each individual
    file contained within, also includes the metadata dict if packaged in the zip"""
    extension = get_file_ext(file_name)
    metadata: dict[str, Any] = {}
    directory_path = os.path.dirname(file_name)

    file_content = get_default_file_store(db_session).read_file(file_name, mode="b")

    if extension == ".zip":
        for file_info, file, metadata in load_files_from_zip(
            file_content, ignore_dirs=True
        ):
            yield os.path.join(directory_path, file_info.filename), file, metadata
    elif extension in [
        ".txt",
        ".md",
        ".mdx",
        ".pdf",
        ".docx",
        ".pptx",
        ".xlsx",
        ".csv",
        ".eml",
        ".epub",
    ]:
        yield file_name, file_content, metadata
    else:
        logger.warning(f"Skipping file '{file_name}' with extension '{extension}'")


def _process_file(
    file_name: str,
    file: IO[Any],
    metadata: dict[str, Any] | None = None,
    pdf_pass: str | None = None,
) -> list[Document]:
    extension = get_file_ext(file_name)
    if not check_file_ext_is_valid(extension):
        logger.warning(f"Skipping file '{file_name}' with extension '{extension}'")
        return []

    file_metadata: dict[str, Any] = {}

    if extension == ".pdf":
        file_content_raw = read_pdf_file(
            file=file, file_name=file_name, pdf_pass=pdf_pass
        )

    elif extension == ".docx":
        file_content_raw = docx2txt.process(file)

    elif extension == ".pptx":
        presentation = pptx.Presentation(file)
        text_content = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            extracted_text = f"\nSlide {slide_number}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    extracted_text += shape.text + "\n"

            text_content.append(extracted_text)
        file_content_raw = "\n\n".join(text_content)

    elif extension == ".xlsx":
        workbook = openpyxl.load_workbook(file)
        text_content = []
        for sheet in workbook.worksheets:
            sheet_string = "\n".join(
                ",".join(map(str, row))
                for row in sheet.iter_rows(min_row=1, values_only=True)
            )
            text_content.append(sheet_string)
        file_content_raw = "\n\n".join(text_content)

    elif extension == ".csv":
        text_file = io.TextIOWrapper(file, encoding=detect_encoding(file))
        reader = csv.reader(text_file)
        file_content_raw = "\n".join([",".join(row) for row in reader])

    elif extension == ".eml":
        text_file = io.TextIOWrapper(file, encoding=detect_encoding(file))
        parser = EmailParser()
        message = parser.parse(text_file)

        text_content = []
        for part in message.walk():
            if part.get_content_type().startswith("text/plain"):
                text_content.append(part.get_payload())
        file_content_raw = "\n\n".join(text_content)

    elif extension == ".epub":
        with zipfile.ZipFile(file) as epub:
            text_content = []
            for item in epub.infolist():
                if item.filename.endswith(".xhtml") or item.filename.endswith(".html"):
                    with epub.open(item) as html_file:
                        soup = BeautifulSoup(html_file, "html.parser")
                        text_content.append(soup.get_text())
            file_content_raw = "\n\n".join(text_content)
    else:
        encoding = detect_encoding(file)
        file_content_raw, file_metadata = read_file(file, encoding=encoding)
    all_metadata = {**metadata, **file_metadata} if metadata else file_metadata

    # If this is set, we will show this in the UI as the "name" of the file
    file_display_name_override = all_metadata.get("file_display_name")

    time_updated = all_metadata.get("time_updated", datetime.now(timezone.utc))
    if isinstance(time_updated, str):
        time_updated = time_str_to_utc(time_updated)

    dt_str = all_metadata.get("doc_updated_at")
    final_time_updated = time_str_to_utc(dt_str) if dt_str else time_updated

    # Metadata tags separate from the Danswer specific fields
    metadata_tags = {
        k: v
        for k, v in all_metadata.items()
        if k
        not in [
            "time_updated",
            "doc_updated_at",
            "link",
            "primary_owners",
            "secondary_owners",
            "filename",
            "file_display_name",
        ]
    }

    p_owner_names = all_metadata.get("primary_owners")
    s_owner_names = all_metadata.get("secondary_owners")
    p_owners = (
        [BasicExpertInfo(display_name=name) for name in p_owner_names]
        if p_owner_names
        else None
    )
    s_owners = (
        [BasicExpertInfo(display_name=name) for name in s_owner_names]
        if s_owner_names
        else None
    )

    return [
        Document(
            id=f"FILE_CONNECTOR__{file_name}",  # add a prefix to avoid conflicts with other connectors
            sections=[
                Section(link=all_metadata.get("link"), text=file_content_raw.strip())
            ],
            source=DocumentSource.FILE,
            semantic_identifier=file_display_name_override
            or os.path.basename(file_name),
            doc_updated_at=final_time_updated,
            primary_owners=p_owners,
            secondary_owners=s_owners,
            # currently metadata just houses tags, other stuff like owners / updated at have dedicated fields
            metadata=metadata_tags,
        )
    ]


class LocalFileConnector(LoadConnector):
    def __init__(
        self,
        file_locations: list[Path | str],
        batch_size: int = INDEX_BATCH_SIZE,
    ) -> None:
        self.file_locations = [Path(file_location) for file_location in file_locations]
        self.batch_size = batch_size
        self.pdf_pass: str | None = None

    def load_credentials(self, credentials: dict[str, Any]) -> dict[str, Any] | None:
        self.pdf_pass = credentials.get("pdf_password")
        return None

    def load_from_state(self) -> GenerateDocumentsOutput:
        documents: list[Document] = []
        with Session(get_sqlalchemy_engine()) as db_session:
            for file_path in self.file_locations:
                current_datetime = datetime.now(timezone.utc)
                files = _read_files_and_metadata(
                    file_name=str(file_path), db_session=db_session
                )

                for file_name, file, metadata in files:
                    metadata["time_updated"] = metadata.get(
                        "time_updated", current_datetime
                    )
                    documents.extend(
                        _process_file(file_name, file, metadata, self.pdf_pass)
                    )

                    if len(documents) >= self.batch_size:
                        yield documents
                        documents = []

            if documents:
                yield documents


if __name__ == "__main__":
    connector = LocalFileConnector(file_locations=[os.environ["TEST_FILE"]])
    connector.load_credentials({"pdf_password": os.environ["PDF_PASSWORD"]})

    document_batches = connector.load_from_state()
    print(next(document_batches))
